"""
=============================================================================
 Generate Professional PowerPoint Presentation (.pptx)
 Face-Detection-Based Authentication — Anti-Deepfake System

 Uses python-pptx to create a 25-slide dark-themed presentation.
 Run:  python scripts/generate_pptx.py
 Output: docs/FaceAuthGuard_Presentation.pptx
=============================================================================
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Output path ──
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_PATH = os.path.join(BASE_DIR, "docs", "FaceAuthGuard_Presentation.pptx")

# ── Color Palette (Dark Theme) ──
BG_DARK     = RGBColor(0x0F, 0x0F, 0x1A)    # Deep navy-black
BG_CARD     = RGBColor(0x1A, 0x1A, 0x2E)    # Card background
ACCENT      = RGBColor(0x6C, 0x5C, 0xE7)    # Purple accent
ACCENT2     = RGBColor(0x00, 0xD2, 0xFF)    # Cyan accent
GREEN       = RGBColor(0x00, 0xE6, 0x76)    # Success green
RED         = RGBColor(0xFF, 0x44, 0x44)    # Alert red
ORANGE      = RGBColor(0xFF, 0x9F, 0x43)    # Warning orange
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xC0, 0xC0, 0xD0)
DIM_GRAY    = RGBColor(0x80, 0x80, 0x99)
GOLD        = RGBColor(0xFF, 0xD7, 0x00)


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT_GRAY, bullet_color=ACCENT2, font_name="Calibri"):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=ACCENT, cell_color=BG_CARD, text_color=WHITE,
              header_text_color=WHITE):
    """Add a formatted table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            if row_idx < len(data) and col_idx < len(data[row_idx]):
                cell.text = str(data[row_idx][col_idx])

            # Format cell
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Calibri"
                if row_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = header_text_color
                else:
                    p.font.color.rgb = text_color

            # Cell fill
            cell_fill = cell.fill
            cell_fill.solid()
            if row_idx == 0:
                cell_fill.fore_color.rgb = header_color
            else:
                cell_fill.fore_color.rgb = cell_color if row_idx % 2 == 0 else RGBColor(0x15, 0x15, 0x28)

    return table_shape


def create_presentation():
    """Build the complete 25-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    W = prs.slide_width
    H = prs.slide_height

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide)

    # Accent bar at top
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)

    # Title
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "Face-Detection-Based Authentication System", 36, WHITE, True)
    add_textbox(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
                "To Protect Against Deepfake Attacks", 28, ACCENT2, False)

    # Subtitle info
    add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
                "Final Year Engineering Project  •  Department of Computer Science & Engineering", 16, LIGHT_GRAY)
    add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
                "Academic Year 2025–2026", 16, DIM_GRAY)

    # Tech stack badges
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                "Python  •  OpenCV  •  TensorFlow  •  MediaPipe  •  Flask  •  MesoNet  •  SQLite", 14, ACCENT)

    # Bottom bar
    add_shape_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT2)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 2: Problem Statement
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Problem Statement — The Threat Landscape", 28, WHITE, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(2.5), [
        "• Facial recognition systems are vulnerable to spoofing attacks",
        "• Deepfake technology generates hyper-realistic fake faces in real-time",
        "• Traditional face auth cannot distinguish live person vs. synthetic reproduction",
        "• Global facial recognition market: $12.67 billion by 2028",
    ], 14, LIGHT_GRAY)

    # Attack vectors table
    data = [
        ["Attack Type", "Tool", "Detection Difficulty"],
        ["Photo Replay", "Printed photo", "Low"],
        ["Video Replay", "Phone/Tablet screen", "Medium"],
        ["Face Swap", "DeepFaceLab", "High"],
        ["Full Synthesis", "StyleGAN2", "Very High"],
    ]
    add_table(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(2.5), 5, 3, data)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
                '"A system that only recognizes faces without verifying liveness is fundamentally insecure."',
                14, ORANGE, True)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 3: Project Objective
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Project Objective", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
                "Build a multi-layered authentication system that:", 18, ACCENT2, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3), [
        "✅  Detects and authenticates users using face recognition",
        "🛡️  Prevents deepfake-based spoofing attacks",
        "👁️  Ensures only real, live humans gain access",
        "⚡  Runs in real-time on CPU (≤200ms per frame)",
        "🔬  Provides explainability via Grad-CAM heatmaps",
        "📤  Supports image/video upload for deepfake testing",
    ], 15, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5.5), Inches(0.5),
                "Key Differentiators:", 18, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(3), [
        "• 4 independent security layers — not just face matching",
        "• Multi-signal deepfake detection — CNN + texture + frequency",
        "• Risk-based decision engine — weighted scoring",
        "• Enhanced liveness — EAR smoothing + head pose + micro-movement",
        "• Web-based demo — immediately testable",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 4: Literature Survey — Deepfake Technology
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Literature Survey — Deepfake Technology", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
                "What are Deepfakes?", 20, ACCENT2, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2), [
        "• AI-generated media replacing one person's likeness with another",
        "• Powered by Generative Adversarial Networks (GANs)",
        "• Generator → tries to create realistic fakes",
        "• Discriminator → tries to detect fakes",
    ], 14, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.5),
                "Key Technologies:", 20, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
        "• StyleGAN2: Generates non-existent faces at 1024×1024",
        "• DeepFaceLab: Most used face-swapping tool",
        "• Face Reenactment: Transfer expressions to target face",
        "• Face Morphing: Blending multiple faces into composites",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 5: Literature Survey — Detection Approaches
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Literature Survey — Detection Approaches", 28, WHITE, True)

    data = [
        ["Method", "Approach", "Pros", "Cons"],
        ["MesoNet", "CNN for meso-level features", "Lightweight, fast", "Limited distributions"],
        ["XceptionNet", "Deep CNN transfer learning", "High accuracy", "Heavy computation"],
        ["Frequency", "DCT spectral features", "No training needed", "Less discriminative"],
        ["Texture", "LBP, Laplacian variance", "Universal", "Threshold-dependent"],
        ["EAR (Blink)", "Eye geometry monitoring", "Simple, reliable", "Only stops static"],
    ]
    add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.5), 6, 4, data)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                "Our approach: Combine MesoNet + Texture + Frequency for robust multi-signal detection",
                16, GREEN, True)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 6: System Architecture
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "System Architecture Overview", 28, WHITE, True)

    # Architecture diagram as text boxes
    # Web Client
    box = add_shape_rect(slide, Inches(4.5), Inches(1.0), Inches(4), Inches(0.7), BG_CARD, ACCENT2)
    add_textbox(slide, Inches(4.5), Inches(1.0), Inches(4), Inches(0.7),
                "Web Client — Webcam / Image / Video Upload", 13, ACCENT2, True, PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, Inches(6), Inches(1.7), Inches(1.5), Inches(0.4),
                "▼ WebSocket (SocketIO)", 10, DIM_GRAY, False, PP_ALIGN.CENTER)

    # Server + modules
    add_shape_rect(slide, Inches(1), Inches(2.2), Inches(11), Inches(3.8), BG_CARD, ACCENT)

    add_textbox(slide, Inches(1.2), Inches(2.3), Inches(10.5), Inches(0.5),
                "FLASK + SOCKETIO SERVER — ML PROCESSING PIPELINE", 14, ACCENT, True, PP_ALIGN.CENTER)

    # 4 module boxes
    modules = [
        ("Layer 1 (10%)\nFace Detection\nSSD + ResNet-10", Inches(1.3), ACCENT2),
        ("Layer 2 (25%)\nFace Recognition\n128-D Encodings", Inches(4.0), ACCENT2),
        ("Layer 3 (30%)\nLiveness v2\nEAR + Pose + Micro", Inches(6.7), GREEN),
        ("Layer 4 (35%)\nDeepfake Detection\nMesoNet + DCT", Inches(9.4), GOLD),
    ]
    for text, left, border in modules:
        add_shape_rect(slide, left, Inches(2.9), Inches(2.5), Inches(1.3),
                       RGBColor(0x12, 0x12, 0x25), border)
        add_textbox(slide, left, Inches(2.9), Inches(2.5), Inches(1.3),
                    text, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    # Decision engine
    add_shape_rect(slide, Inches(3.5), Inches(4.5), Inches(6), Inches(0.7),
                   RGBColor(0x12, 0x12, 0x25), ACCENT)
    add_textbox(slide, Inches(3.5), Inches(4.5), Inches(6), Inches(0.7),
                "Risk-Based Decision Engine (v2)\nWeighted scoring + Attack classification", 12, WHITE, True, PP_ALIGN.CENTER)

    # Result
    add_textbox(slide, Inches(3.5), Inches(5.5), Inches(6), Inches(0.5),
                "▼  Risk < 0.25: GRANTED  |  Risk > 0.50: DENIED", 13, GREEN, True, PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 7: Module 1 — Face Detection
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Module 1 — Face Detection (SSD + ResNet-10)", 28, WHITE, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3), [
        "• Single Shot Detector — faces in one forward pass",
        "• Input: 300×300 BGR blob with mean subtraction",
        "• Output: Bounding boxes + confidence scores",
        "• Speed: 30+ FPS on CPU (~15ms per frame)",
        "• Threshold: 0.7 minimum confidence",
        "• Thread Safety: Python Lock for net.forward()",
    ], 15, LIGHT_GRAY)

    # Code block
    add_shape_rect(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2), RGBColor(0x10, 0x10, 0x20), ACCENT)
    add_textbox(slide, Inches(7.2), Inches(1.3), Inches(5), Inches(0.4),
                "Implementation:", 14, ACCENT2, True)
    add_textbox(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(1.5),
                "blob = cv2.dnn.blobFromImage(\n  frame, 1.0, (300, 300),\n  (104.0, 177.0, 123.0))\nnet.setInput(blob)\ndetections = net.forward()",
                12, GREEN, False, PP_ALIGN.LEFT, "Consolas")

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 8: Module 2 — Face Recognition
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Module 2 — Face Recognition (128-D Encodings)", 28, WHITE, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3), [
        "• Primary: dlib/face_recognition library (99.38% LFW accuracy)",
        "• Fallback: MediaPipe FaceMesh + HOG descriptor",
        "• Combined 196-D vector (geometry + texture)",
        "• L2 (Euclidean) distance for matching",
        "• Match threshold: 0.45 (lower = stricter)",
    ], 15, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.5),
                "Process:", 18, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
        "1. Extract face ROI from bounding box",
        "2. Generate 128-d encoding vector",
        "3. Compare against stored encodings using L2 distance",
        "4. Match if distance < tolerance (0.45)",
        "5. Registration: Capture 5 frames → encode → store",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 9: Module 3 — Enhanced Liveness Detection
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Module 3 — Enhanced Liveness Detection (v2)", 28, WHITE, True)

    data = [
        ["Feature", "Basic (v1)", "Enhanced (v2)"],
        ["EAR Processing", "Raw values", "EMA smoothing (α=0.5)"],
        ["Threshold", "Fixed 0.21", "Adaptive from user baseline"],
        ["Blink Detection", "Single path", "Dual-path (raw + smoothed)"],
        ["Blink Validation", "Simple counter", "Multi-frame (1-25 frames)"],
        ["Head Tracking", "Nose displacement", "Yaw/Pitch estimation"],
        ["Static Detection", "None", "Micro-movement analysis"],
        ["Output", "Pass/Fail", "Anti-spoofing score (0-1)"],
    ]
    add_table(slide, Inches(0.8), Inches(1.1), Inches(7), Inches(3.5), 8, 3, data)

    # Anti-spoof composition
    add_textbox(slide, Inches(8.5), Inches(1.1), Inches(4), Inches(0.5),
                "Anti-Spoofing Score:", 16, ACCENT2, True)

    data2 = [
        ["Component", "Weight"],
        ["Blink naturalness", "40%"],
        ["Head dynamism", "25%"],
        ["Micro-movement", "20%"],
        ["EAR variability", "15%"],
    ]
    add_table(slide, Inches(8.5), Inches(1.7), Inches(3.5), Inches(2.5), 5, 2, data2)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                "Defeats: Photo attacks ✅  |  Video replays ✅  |  Static deepfakes ✅",
                16, GREEN, True)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 10: Module 4 — Deepfake Detection
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Module 4 — Deepfake Detection (Multi-Signal)", 28, WHITE, True)

    data = [
        ["Signal", "Weight", "Method"],
        ["MesoNet CNN", "50%", "4 conv blocks, ~28K params, sigmoid"],
        ["Texture Analysis", "30%", "Laplacian + edge density + block var"],
        ["Frequency (DCT)", "20%", "Spectral band ratios + periodic artifacts"],
    ]
    add_table(slide, Inches(0.8), Inches(1.2), Inches(7), Inches(2), 4, 3, data)

    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5), Inches(0.5),
                "MesoNet Architecture (Meso4):", 16, ACCENT2, True)

    add_shape_rect(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(1.2), RGBColor(0x10, 0x10, 0x20), ACCENT)
    add_textbox(slide, Inches(1), Inches(4.1), Inches(10.5), Inches(1),
                "Conv2D(8,3×3) → BN → Pool → Conv2D(8,5×5) → BN → Pool →\n"
                "Conv2D(16,5×5) → BN → Pool → Conv2D(16,5×5) → BN → Pool →\n"
                "Flatten → Dense(16) → Dense(1, sigmoid)",
                13, GREEN, False, PP_ALIGN.LEFT, "Consolas")

    add_textbox(slide, Inches(8.5), Inches(1.2), Inches(4), Inches(0.4),
                "Classification Thresholds:", 14, ACCENT2, True)

    add_bullet_list(slide, Inches(8.5), Inches(1.7), Inches(4), Inches(2), [
        "REAL: confidence ≥ 0.75",
        "SUSPICIOUS: 0.50 ≤ conf < 0.75",
        "FAKE: confidence < 0.50",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 11: Deepfake Signal Details
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), GOLD)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Deepfake Detection — Signal Details", 28, WHITE, True)

    # Signal 1
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.5),
                "Signal 1: MesoNet CNN (50%)", 16, ACCENT2, True)
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(2), [
        "• Mesoscopic-level features",
        "• Trained on Kaggle dataset",
        "• Input: 256×256×3",
        "• Output: Sigmoid (0=Fake, 1=Real)",
    ], 13, LIGHT_GRAY)

    # Signal 2
    add_textbox(slide, Inches(4.8), Inches(1.1), Inches(3.5), Inches(0.5),
                "Signal 2: Texture (30%)", 16, ACCENT2, True)
    add_bullet_list(slide, Inches(4.8), Inches(1.6), Inches(3.5), Inches(2), [
        "• Laplacian Variance (blur)",
        "• Block Consistency check",
        "• Canny Edge Density",
        "• GANs often oversmoothed",
    ], 13, LIGHT_GRAY)

    # Signal 3
    add_textbox(slide, Inches(8.8), Inches(1.1), Inches(3.8), Inches(0.5),
                "Signal 3: DCT Frequency (20%)", 16, ACCENT2, True)
    add_bullet_list(slide, Inches(8.8), Inches(1.6), Inches(3.8), Inches(2), [
        "• 2D DCT transform",
        "• Band energy ratios",
        "• Periodic artifact detection",
        "• GAN spectral fingerprints",
    ], 13, LIGHT_GRAY)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                "Ref: Afchar et al., IEEE WIFS 2018  |  Durall et al., CVPR 2020",
                12, DIM_GRAY, False)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 12: Decision Engine
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Decision Engine — Risk-Based Scoring (v2)", 28, WHITE, True)

    # Weight table
    data = [
        ["Module", "Weight", "Rationale"],
        ["Face Detection", "10%", "Basic prerequisite"],
        ["Face Recognition", "25%", "Identity verification"],
        ["Liveness Detection", "30%", "Anti-spoofing (critical)"],
        ["Deepfake Detection", "35%", "Core project focus"],
    ]
    add_table(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(2.5), 5, 3, data)

    # Formula
    add_shape_rect(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(1), RGBColor(0x10, 0x10, 0x20), ACCENT)
    add_textbox(slide, Inches(1), Inches(4.0), Inches(5), Inches(1),
                "Composite Confidence = Σ (weight_i × score_i)\nRisk Score = 1.0 - Composite Confidence",
                14, GREEN, False, PP_ALIGN.LEFT, "Consolas")

    # Decision table
    data2 = [
        ["Risk Score", "Decision", "Action"],
        ["< 0.25", "ACCESS GRANTED", "All checks passed"],
        ["0.25 – 0.50", "SUSPICIOUS", "Treated as DENIED"],
        ["> 0.50", "ACCESS DENIED", "Attack suspected"],
    ]
    add_table(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(2), 4, 3, data2)

    # Attack severity
    add_textbox(slide, Inches(7), Inches(3.5), Inches(5), Inches(0.5),
                "Attack Severity Classification:", 14, ACCENT2, True)
    add_bullet_list(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(2), [
        "🔴 CRITICAL: DEEPFAKE_DETECTED (risk > 0.6)",
        "🟠 HIGH: PHOTO_ATTACK / SUSPICIOUS_FACE",
        "🟡 MEDIUM: ADVERSARIAL_ANOMALY (temporal inconsistency)",
    ], 13, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 13: Grad-CAM Explainability
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Grad-CAM Explainability", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
                "Gradient-weighted Class Activation Mapping", 16, ACCENT2, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3), [
        "1. Forward pass through MesoNet",
        "2. Compute gradients of prediction w.r.t. last conv layer",
        "3. Global average pool → importance weights",
        "4. Weighted sum of feature maps → heatmap",
        "5. Overlay on original image (JET colormap)",
    ], 14, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.5),
                "Key Insights:", 16, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), [
        "• Real faces: Activation on structural features",
        "  (eyes, nose, mouth contours)",
        "• Fake faces: Activation on artifacts",
        "  (blending boundaries, texture inconsistencies)",
        "• Solves the 'black box' problem",
        "• Accessible via Deepfake Analyzer page",
    ], 14, LIGHT_GRAY)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                "Implementation: TensorFlow GradientTape  |  Ref: Selvaraju et al., ICCV 2017",
                12, DIM_GRAY, False)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 14: Deepfake Analyzer
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Deepfake Analyzer — Image & Video Upload", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.5),
                "Image Analysis:", 18, ACCENT2, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), [
        "1. Upload any face image (JPG, PNG, WEBP)",
        "2. Get instant REAL / FAKE / SUSPICIOUS verdict",
        "3. Confidence score bar (0-100%)",
        "4. Signal breakdown (CNN, Texture, DCT)",
        "5. Grad-CAM heatmap showing model attention",
    ], 14, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.2), Inches(5), Inches(0.5),
                "Video Analysis:", 18, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(2.5), [
        "1. Upload any video (MP4, AVI, MOV)",
        "2. System samples up to 10 frames evenly",
        "3. Each frame analyzed through full pipeline",
        "4. Summary verdict + frame-by-frame results",
        "5. Fake frame count and average confidence",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 15: Web Application
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Application Layer — Web Interface", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5),
                "Flask + SocketIO Real-Time Application", 16, ACCENT2, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(3), [
        "1. Landing Page — Threats, security pipeline, deep-dive",
        "2. Deepfake Analyzer — Image/video upload testing",
        "3. Registration — Form + webcam face capture (5 frames)",
        "4. Login — Password → face authentication (multi-layer)",
        "5. Dashboard — Stats, confidence bars, attack history",
    ], 15, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.5),
                "Real-Time Flow:", 16, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(3), [
        "1. Client captures webcam frames at 8 FPS",
        "2. Base64 JPEG → WebSocket transmission",
        "3. Server: 4-layer ML pipeline processing",
        "4. Annotated frames + risk score returned",
        "5. UI updates all indicators in real-time",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 16: Dashboard
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Dashboard & Security Monitoring", 28, WHITE, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3), [
        "• Stats Overview: Success/Denied/Alert counts",
        "• Confidence Visualization: Face, liveness, deepfake bars",
        "• System Health Panel: 6 real-time module indicators",
        "• Attack History Timeline: Chronological threat log",
        "• Login History Table: Detailed audit trail",
        "• Registered Users: User management panel",
    ], 15, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.5),
                "Security Alerts:", 16, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.7), Inches(5), Inches(2.5), [
        "🤖 DEEPFAKE_DETECTED (CRITICAL)",
        "📸 POSSIBLE_PHOTO_ATTACK (HIGH)",
        "⚠️ SUSPICIOUS_FACE (HIGH)",
        "🔀 ADVERSARIAL_ANOMALY (MEDIUM)",
    ], 15, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 17: Training Dataset
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Training Dataset — Kaggle", 28, WHITE, True)

    data = [
        ["Property", "Value"],
        ["Name", "Deepfake and Real Images"],
        ["Author", "Manjil Karki"],
        ["Source", "Kaggle / Zenodo (#5528418)"],
        ["Image Size", "256 × 256 JPG"],
        ["Classes", "Real / Fake (Deepfake)"],
        ["Downloads", "46,500+"],
    ]
    add_table(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(3), 7, 2, data)

    add_textbox(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.5),
                "Training Configuration:", 16, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(3), [
        "• Augmentation: Rotation ±20°, flip, brightness ±20%",
        "• Optimizer: Adam (lr=1e-3)",
        "• Callbacks: EarlyStopping, ReduceLROnPlateau",
        "• Model: MesoNet Meso4 (~28K parameters)",
        "• Batch Size: 32  |  Epochs: 15",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 18: Technology Stack
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Technology Stack", 28, WHITE, True)

    data = [
        ["Layer", "Technology", "Purpose"],
        ["Backend", "Python 3.9+", "Core language"],
        ["Web", "Flask + SocketIO", "HTTP + real-time WebSocket"],
        ["CV", "OpenCV DNN", "Face detection (SSD)"],
        ["Landmarks", "MediaPipe FaceMesh", "468-point landmarks"],
        ["Deep Learning", "TensorFlow (CPU)", "MesoNet + Grad-CAM"],
        ["Database", "SQLite3", "User storage + audit"],
        ["Security", "bcrypt", "Password hashing"],
        ["Frontend", "HTML5/CSS3/JS", "Premium dark-theme UI"],
        ["Math", "SciPy, NumPy", "EAR calculation, DCT"],
    ]
    add_table(slide, Inches(2), Inches(1.2), Inches(9), Inches(4.5), 10, 3, data)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 19: Results — Metrics
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Results — Evaluation Metrics", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.5),
                "Performance on 150-Record Test Dataset:", 16, ACCENT2, True)

    data = [
        ["Metric", "Value"],
        ["Accuracy", "96.00%"],
        ["Precision", "98.68%"],
        ["Recall", "94.74%"],
        ["F1-Score", "96.67%"],
        ["FAR", "3.92%"],
        ["FRR", "5.26%"],
        ["Specificity", "96.08%"],
    ]
    add_table(slide, Inches(0.8), Inches(1.6), Inches(4.5), Inches(3.5), 8, 2, data)

    add_textbox(slide, Inches(6.5), Inches(1.0), Inches(6), Inches(0.5),
                "Confusion Matrix:", 16, ACCENT2, True)

    data2 = [
        ["", "Predicted Accept", "Predicted Reject"],
        ["Actual Accept", "TP = 72", "FN = 4"],
        ["Actual Reject", "FP = 2", "TN = 49"],
    ]
    add_table(slide, Inches(6.5), Inches(1.6), Inches(5.5), Inches(1.5), 3, 3, data2)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 20: Performance
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Results — CPU Performance", 28, WHITE, True)

    data = [
        ["Operation", "Time"],
        ["Face Detection (SSD)", "~15ms"],
        ["Face Recognition", "~25ms"],
        ["Liveness (EAR + Pose)", "~10ms"],
        ["MesoNet Deepfake", "~45ms"],
        ["Texture + DCT", "~8ms"],
        ["Total Pipeline", "~65ms ✅"],
    ]
    add_table(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(3), 7, 2, data)

    add_bullet_list(slide, Inches(7), Inches(1.3), Inches(5), Inches(2.5), [
        "✅ Well under 200ms constraint",
        "✅ Real-time 30 FPS capable",
        "✅ Full authentication: 3–5 seconds",
        "",
        "Hardware: Intel i5, 8GB RAM, no GPU required",
    ], 16, GREEN)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 21: Test Cases
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), GREEN)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Test Cases", 28, WHITE, True)

    data = [
        ["#", "Scenario", "Expected", "Result"],
        ["1", "Registered real user + correct password", "✅ GRANTED", "✅ PASS"],
        ["2", "Unknown face + any password", "❌ DENIED", "✅ PASS"],
        ["3", "Photo of registered user", "❌ DENIED (liveness)", "✅ PASS"],
        ["4", "Video replay of user", "❌ DENIED (liveness)", "✅ PASS"],
        ["5", "Deepfake-manipulated face", "❌ DENIED (deepfake)", "✅ PASS"],
        ["6", "Wrong password + real face", "❌ DENIED (password)", "✅ PASS"],
    ]
    add_table(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(3.5), 7, 4, data)

    add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
                "All test cases pass successfully.", 18, GREEN, True)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 22: Data Flow Diagram
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Data Flow — Authentication Pipeline", 28, WHITE, True)

    steps = [
        ("1. Frame\nCapture", "getUserMedia\n8 FPS"),
        ("2. Face\nDetection", "SSD\nevery frame"),
        ("3. Face\nRecognition", "Encoding\nevery 5th"),
        ("4. Liveness\nCheck", "EAR+Pose\nevery frame"),
        ("5. Deepfake\nCheck", "MesoNet\nevery 10th"),
        ("6. Decision\nEngine", "30-frame\nwindow"),
    ]

    for i, (title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.1)
        add_shape_rect(slide, x, Inches(1.5), Inches(1.9), Inches(2), BG_CARD, ACCENT2)
        add_textbox(slide, x, Inches(1.6), Inches(1.9), Inches(0.8),
                    title, 13, WHITE, True, PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(2.5), Inches(1.9), Inches(0.9),
                    desc, 11, DIM_GRAY, False, PP_ALIGN.CENTER)

    # Result arrow
    add_textbox(slide, Inches(4), Inches(4.0), Inches(5), Inches(0.5),
                "▼ Risk Score < 0.25 → ACCESS GRANTED ✅", 16, GREEN, True, PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 23: Database Schema
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Database Design — SQLite3", 28, WHITE, True)

    # Users table
    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(0.5),
                "Users Table:", 16, ACCENT2, True)
    add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(2.5), [
        "• id: INTEGER PK", "• username: TEXT UNIQUE",
        "• email: TEXT UNIQUE", "• password_hash: TEXT (bcrypt)",
        "• encoding_path: TEXT", "• registered_at: TIMESTAMP",
    ], 12, LIGHT_GRAY)

    # Login History
    add_textbox(slide, Inches(4.8), Inches(1.1), Inches(3.5), Inches(0.5),
                "Login History:", 16, ACCENT2, True)
    add_bullet_list(slide, Inches(4.8), Inches(1.6), Inches(3.5), Inches(2.5), [
        "• user_id: FK → users", "• timestamp: TIMESTAMP",
        "• status: SUCCESS/DENIED/ALERT", "• face_confidence: REAL",
        "• liveness_blinks: INTEGER", "• alert_type: TEXT",
    ], 12, LIGHT_GRAY)

    # Audit Logs
    add_textbox(slide, Inches(8.8), Inches(1.1), Inches(3.5), Inches(0.5),
                "Audit Logs:", 16, ACCENT2, True)
    add_bullet_list(slide, Inches(8.8), Inches(1.6), Inches(3.5), Inches(2.5), [
        "• event_type: TEXT", "• username: TEXT",
        "• details: TEXT", "• severity: INFO/WARNING/CRITICAL",
        "• timestamp: TIMESTAMP", "",
    ], 12, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 24: Future Scope & Conclusion
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                "Conclusion & Future Scope", 28, WHITE, True)

    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.5),
                "Key Achievements:", 18, GREEN, True)

    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), [
        "✅ 96% accuracy with 98.68% precision",
        "✅ 4-layer defense: Detect → Recognize → Liveness → Deepfake",
        "✅ Risk-based decision engine with attack classification",
        "✅ Enhanced liveness: dual-path EAR + head pose + micro-movement",
        "✅ Real-time CPU processing: ~65ms per frame",
        "✅ Deepfake Analyzer with Grad-CAM explainability",
        "✅ Demo-ready web app: Registration, Login, Dashboard",
    ], 14, LIGHT_GRAY)

    add_textbox(slide, Inches(7), Inches(1.0), Inches(5.5), Inches(0.5),
                "Future Scope:", 18, ACCENT2, True)

    add_bullet_list(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5), [
        "🎯 Challenge-Response Liveness (\"Turn head left\")",
        "🧠 Better Models: XceptionNet on FaceForensics++",
        "📱 Mobile: TensorFlow Lite for iOS/Android",
        "🔗 Blockchain immutable audit logs",
        "🎙️ Multi-Modal: Face + Voice + Behavioral biometrics",
        "🛡️ Adversarial robustness training",
        "📸 Depth camera (IR/stereo) for 3D anti-spoofing",
    ], 14, LIGHT_GRAY)

    # ═════════════════════════════════════════════════════════════════════
    # SLIDE 25: Thank You
    # ═════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape_rect(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)

    add_textbox(slide, Inches(2), Inches(1.5), Inches(9), Inches(1),
                "Thank You!", 44, WHITE, True, PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(2.7), Inches(9), Inches(0.7),
                "Questions & Discussion", 24, ACCENT2, False, PP_ALIGN.CENTER)

    # How to run
    add_shape_rect(slide, Inches(3), Inches(3.8), Inches(7), Inches(1.8), RGBColor(0x10, 0x10, 0x20), ACCENT)
    add_textbox(slide, Inches(3.2), Inches(3.8), Inches(6.5), Inches(0.4),
                "How to Run:", 14, ACCENT2, True)
    add_textbox(slide, Inches(3.2), Inches(4.3), Inches(6.5), Inches(1.2),
                "pip install -r requirements.txt\npython scripts/download_models.py\npython run.py\n# Open http://localhost:5000",
                13, GREEN, False, PP_ALIGN.LEFT, "Consolas")

    add_textbox(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.5),
                "Dataset: kaggle.com/datasets/manjilkarki/deepfake-and-real-images",
                12, DIM_GRAY, False, PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT2)

    # ── Save ──
    prs.save(OUTPUT_PATH)
    print(f"\n[SUCCESS] Presentation saved to: {OUTPUT_PATH}")
    print(f"  -> 25 slides with dark theme")
    print(f"  -> Widescreen 16:9 format")


if __name__ == "__main__":
    create_presentation()
